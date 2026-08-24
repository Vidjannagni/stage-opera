#!/usr/bin/env python3
"""Génère la présentation de soutenance au format PowerPoint.

Le fichier .pptx produit est entièrement modifiable dans PowerPoint : chaque
élément est une forme native (zone de texte, rectangle, tableau, image), aucune
diapositive n'est une image aplatie. Relancer ce script écrase le fichier.

    python3 construire_presentation.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ICI = Path(__file__).parent
IMG = ICI / "img"
SORTIE = ICI / "Soutenance_stage_HOUNKPETOHOU.pptx"

# ══════════════════════════════════════════════════════════════════════════════
#  CHARTE — tout le style tient dans ce bloc
# ══════════════════════════════════════════════════════════════════════════════
NAVY    = RGBColor(0x0C, 0x19, 0x30)   # fonds pleins, titres
BLEU    = RGBColor(0x00, 0x52, 0x91)   # couleur structurante (École Centrale)
TEAL    = RGBColor(0x00, 0xBC, 0xAA)   # accent vif
OR      = RGBColor(0xF9, 0xCC, 0x28)   # filet et étiquettes (logo du cabinet)
ARDOISE = RGBColor(0x56, 0x65, 0x7C)   # texte secondaire
BRUME   = RGBColor(0xF2, 0xF5, 0xF9)   # fond des cartes
TRAIT   = RGBColor(0xD9, 0xE0, 0xEA)   # filets discrets
BLANC   = RGBColor(0xFF, 0xFF, 0xFF)
ROUGE   = RGBColor(0xB0, 0x00, 0x20)   # informations restant à compléter

TITRE_POLICE = "Georgia"     # titres : serif éditorial
TEXTE_POLICE = "Calibri"     # texte courant

NB_DIAPOS = 16               # total affiché en pied de page (vérifié à la fin)

L, H = 13.333, 7.5           # diapositive 16:9
MARGE = 0.72
BANDE_H = 1.16               # hauteur du bandeau de titre
HAUT = 1.62                  # première ligne de contenu
BAS = 6.82                   # dernière ligne de contenu
LARG = L - 2 * MARGE         # largeur utile


# ══════════════════════════════════════════════════════════════════════════════
#  Primitives de dessin
# ══════════════════════════════════════════════════════════════════════════════
def rect(diapo, x, y, cx, cy, fond=None, contour=None, epaisseur=0.75,
         forme=MSO_SHAPE.RECTANGLE):
    f = diapo.shapes.add_shape(forme, Inches(x), Inches(y), Inches(cx), Inches(cy))
    f.shadow.inherit = False
    _sans_effet_de_theme(f)
    if fond is None:
        f.fill.background()
    else:
        f.fill.solid()
        f.fill.fore_color.rgb = fond
    if contour is None:
        f.line.fill.background()
    else:
        f.line.color.rgb = contour
        f.line.width = Pt(epaisseur)
    f.text_frame.word_wrap = True
    return f


def _sans_effet_de_theme(forme):
    """Neutralise l'ombre portée héritée du thème (effectRef du style de forme).

    Sans cela, PowerPoint applique l'effet n° 2 du thème : un rendu « années
    2010 » dont on ne veut pas ici, les blocs devant rester parfaitement plats.
    """
    from pptx.oxml.ns import qn
    style = forme._element.find(qn("p:style"))
    if style is None:
        return
    ref = style.find(qn("a:effectRef"))
    if ref is not None:
        ref.set("idx", "0")


def zone(diapo, x, y, cx, cy, ancrage=MSO_ANCHOR.TOP):
    z = diapo.shapes.add_textbox(Inches(x), Inches(y), Inches(cx), Inches(cy))
    tf = z.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ancrage
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def ecrire(tf, morceaux, taille=13, couleur=NAVY, police=TEXTE_POLICE,
           gras=False, italique=False, interligne=1.16, espace_avant=0,
           espace_apres=0, alignement=PP_ALIGN.LEFT, premier=False):
    """Ajoute un paragraphe. `morceaux` : str, ou liste de (texte, options)."""
    p = tf.paragraphs[0] if premier else tf.add_paragraph()
    p.alignment = alignement
    p.line_spacing = interligne
    p.space_before = Pt(espace_avant)
    p.space_after = Pt(espace_apres)
    if isinstance(morceaux, str):
        morceaux = [(morceaux, {})]
    for texte, opt in morceaux:
        r = p.add_run()
        r.text = texte
        f = r.font
        f.name = opt.get("police", police)
        f.size = Pt(opt.get("taille", taille))
        f.bold = opt.get("gras", gras)
        f.italic = opt.get("italique", italique)
        f.color.rgb = opt.get("couleur", couleur)
    return p


def etiquette(tf, texte, couleur=OR, taille=9.5, premier=False, espace_apres=4):
    """Sur-titre en capitales espacées."""
    return ecrire(tf, " ".join(texte.upper()), taille=taille, couleur=couleur,
                  gras=True, espace_apres=espace_apres, premier=premier)


def puces(tf, elements, taille=12.5, couleur=NAVY, puce_couleur=TEAL,
          interligne=1.14, espace=7, premier=False):
    """Liste à puces carrées, avec gras sur le segment avant le tiret cadratin."""
    for i, texte in enumerate(elements):
        if "—" in texte:
            tete, reste = texte.split("—", 1)
            morceaux = [("▪  ", {"couleur": puce_couleur, "gras": True}),
                        (tete.strip() + " ", {"gras": True}),
                        ("— " + reste.strip(), {})]
        else:
            morceaux = [("▪  ", {"couleur": puce_couleur, "gras": True}),
                        (texte, {})]
        ecrire(tf, morceaux, taille=taille, couleur=couleur,
               interligne=interligne, espace_apres=espace,
               premier=(premier and i == 0))


def image_ajustee(diapo, chemin, x, y, cx, cy, cadre=True):
    """Insère une image centrée dans la boîte (x, y, cx, cy), sans déformation."""
    from PIL import Image
    with Image.open(chemin) as im:
        ratio = im.width / im.height
    if cx / cy > ratio:            # boîte plus large que l'image
        h, w = cy, cy * ratio
    else:
        w, h = cx, cx / ratio
    px, py = x + (cx - w) / 2, y + (cy - h) / 2
    img = diapo.shapes.add_picture(str(chemin), Inches(px), Inches(py),
                                   Inches(w), Inches(h))
    if cadre:
        img.line.color.rgb = TRAIT
        img.line.width = Pt(0.75)
    return img


# ══════════════════════════════════════════════════════════════════════════════
#  Gabarits de diapositive
# ══════════════════════════════════════════════════════════════════════════════
def diapo_vierge(prez):
    return prez.slides.add_slide(prez.slide_layouts[6])


def bandeau(diapo, numero, titre, sous_titre=None):
    """Bandeau de titre : pastille numérotée, titre, sous-titre, filet doré."""
    rect(diapo, 0, 0, L, BANDE_H, fond=NAVY)
    rect(diapo, 0, BANDE_H, L, 0.055, fond=OR)

    pastille = rect(diapo, MARGE, 0.30, 0.56, 0.56, fond=None, contour=OR,
                    epaisseur=1.1, forme=MSO_SHAPE.OVAL)
    tf = pastille.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Sans cela, un numéro à deux chiffres se replie sur deux lignes dans le
    # cercle : les marges internes par défaut ne lui laissent pas la place.
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    ecrire(tf, f"{numero:02d}", taille=13, couleur=OR, police=TITRE_POLICE,
           gras=True, alignement=PP_ALIGN.CENTER, premier=True)

    tf = zone(diapo, MARGE + 0.82, 0.24, LARG - 0.82, 0.80, MSO_ANCHOR.MIDDLE)
    ecrire(tf, titre, taille=23, couleur=BLANC, police=TITRE_POLICE, gras=True,
           interligne=1.0, premier=True)
    if sous_titre:
        ecrire(tf, sous_titre, taille=11.5, couleur=RGBColor(0xA9, 0xB8, 0xCC),
               espace_avant=3, interligne=1.0)


def pied(diapo, numero):
    """Filet et mentions de bas de page."""
    rect(diapo, MARGE, H - 0.62, LARG, 0.012, fond=TRAIT)
    tf = zone(diapo, MARGE, H - 0.52, LARG - 0.6, 0.3)
    ecrire(tf, "HOUNKPETOHOU Marius  ·  Soutenance de stage opérateur  ·  "
               "Choubel Consulting", taille=9, couleur=ARDOISE, premier=True)
    tf = zone(diapo, L - MARGE - 0.6, H - 0.52, 0.6, 0.3)
    ecrire(tf, f"{numero} / {NB_DIAPOS}", taille=9, couleur=BLEU, gras=True,
           alignement=PP_ALIGN.RIGHT, premier=True)


def diapo_standard(prez, numero, titre, sous_titre=None):
    d = diapo_vierge(prez)
    bandeau(d, numero, titre, sous_titre)
    pied(d, numero)
    return d


def diapo_photo(prez, photo, opacite_voile=0.82):
    """Diapositive pleine page sur photographie assombrie."""
    d = diapo_vierge(prez)
    from PIL import Image
    with Image.open(photo) as im:
        ratio = im.width / im.height
    if L / H > ratio:
        w, h = L, L / ratio
    else:
        h, w = H, H * ratio
    d.shapes.add_picture(str(photo), Inches((L - w) / 2), Inches((H - h) / 2),
                         Inches(w), Inches(h))
    voile = rect(d, 0, 0, L, H, fond=NAVY)
    voile.fill.transparency = 1 - opacite_voile
    # `transparency` n'est pas exposé par python-pptx : appliqué en XML.
    _transparence(voile, 1 - opacite_voile)
    return d


def _transparence(forme, valeur):
    """Applique une transparence (0 = opaque, 1 = invisible) au remplissage."""
    from pptx.oxml.ns import qn
    solide = forme.fill._xPr.find(qn("a:solidFill"))
    couleur = solide.find(qn("a:srgbClr"))
    alpha = couleur.makeelement(qn("a:alpha"), {"val": str(int((1 - valeur) * 100000))})
    couleur.append(alpha)


def carte(diapo, x, y, cx, cy, titre=None, fond=BRUME, contour=TRAIT,
          bord_gauche=None):
    """Bloc de contenu : fond clair, filet, liseré de couleur optionnel."""
    c = rect(diapo, x, y, cx, cy, fond=fond, contour=contour)
    if bord_gauche is not None:
        rect(diapo, x, y, 0.055, cy, fond=bord_gauche)
    if titre:
        tf = zone(diapo, x + 0.26, y + 0.20, cx - 0.5, 0.3)
        ecrire(tf, titre, taille=12.5, couleur=BLEU, gras=True, premier=True)
    return c


def tableau(diapo, x, y, cx, lignes, largeurs, hauteur_ligne=0.34,
            hauteur_entete=0.40, taille=11):
    """Tableau natif PowerPoint. `lignes[0]` sert d'en-tête."""
    n, m = len(lignes), len(lignes[0])
    forme = diapo.shapes.add_table(n, m, Inches(x), Inches(y), Inches(cx),
                                   Inches(hauteur_entete + (n - 1) * hauteur_ligne))
    t = forme.table
    t.first_row = True
    t.horz_banding = False
    for j, part in enumerate(largeurs):
        t.columns[j].width = Emu(int(Inches(cx) * part / sum(largeurs)))
    for i, ligne in enumerate(lignes):
        t.rows[i].height = Inches(hauteur_entete if i == 0 else hauteur_ligne)
        for j, valeur in enumerate(ligne):
            cel = t.cell(i, j)
            cel.margin_left = Inches(0.12)
            cel.margin_right = Inches(0.10)
            cel.margin_top = cel.margin_bottom = Inches(0.045)
            cel.vertical_anchor = MSO_ANCHOR.MIDDLE
            cel.fill.solid()
            if i == 0:
                cel.fill.fore_color.rgb = NAVY
            else:
                cel.fill.fore_color.rgb = BLANC if i % 2 else BRUME
            tf = cel.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = 1.05
            gras_ligne = valeur.startswith("**")
            r = p.add_run()
            r.text = valeur.replace("**", "")
            r.font.name = TEXTE_POLICE
            r.font.size = Pt(taille)
            r.font.bold = i == 0 or gras_ligne
            r.font.color.rgb = BLANC if i == 0 else NAVY
    return t


def chevrons(diapo, x, y, cx, cy, etapes, actif=None):
    """Suite d'étapes numérotées reliées par des chevrons."""
    n = len(etapes)
    espace = 0.10
    larg = (cx - (n - 1) * espace) / n
    for i, (num, texte) in enumerate(etapes):
        px = x + i * (larg + espace)
        forme = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        f = rect(diapo, px, y, larg, cy,
                 fond=BLEU if actif == i else BRUME,
                 contour=None if actif == i else TRAIT, forme=forme)
        tf = f.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.05)
        ecrire(tf, num, taille=10.5, couleur=BLANC if actif == i else TEAL,
               gras=True, alignement=PP_ALIGN.CENTER, interligne=1.0, premier=True)
        ecrire(tf, texte, taille=10, couleur=BLANC if actif == i else NAVY,
               alignement=PP_ALIGN.CENTER, interligne=1.0, espace_avant=1)


def chiffre_cle(diapo, x, y, cx, valeur, libelle, couleur=BLEU):
    """Grand nombre avec sa légende."""
    tf = zone(diapo, x, y, cx, 0.52)
    ecrire(tf, valeur, taille=27, couleur=couleur, police=TITRE_POLICE,
           gras=True, interligne=1.0, premier=True)
    tf = zone(diapo, x, y + 0.50, cx, 0.4)
    ecrire(tf, libelle, taille=10, couleur=ARDOISE, interligne=1.05, premier=True)


def a_completer(tf, texte, premier=False):
    ecrire(tf, f"[à compléter : {texte}]", taille=10, couleur=ROUGE, gras=True,
           espace_avant=6, premier=premier)


# ══════════════════════════════════════════════════════════════════════════════
#  LES DIAPOSITIVES
# ══════════════════════════════════════════════════════════════════════════════
def diapo_01_titre(prez):
    d = diapo_photo(prez, IMG / "photos/casablanca_finance_city.jpg", 0.86)

    for logo, ancre_droite in ((ICI.parent / "logo.jpg", False),
                               (ICI.parent / "app/static/img/logo-choubel.jpg", True)):
        if not logo.exists():
            continue
        x = L - MARGE - 1.05 if ancre_droite else MARGE
        cadre = rect(d, x - 0.09, 0.42, 1.23, 1.23, fond=BLANC)
        image_ajustee(d, logo, x, 0.51, 1.05, 1.05, cadre=False)

    tf = zone(d, MARGE, 2.55, 10.4, 0.4)
    etiquette(tf, "Soutenance de stage opérateur — première année", premier=True)

    tf = zone(d, MARGE, 3.02, 11.0, 1.9)
    ecrire(tf, "Concevoir un outil d'aide à la décision", taille=36,
           couleur=BLANC, police=TITRE_POLICE, gras=True, interligne=1.08,
           premier=True)
    ecrire(tf, "pour un cabinet de conseil immobilier", taille=36,
           couleur=BLANC, police=TITRE_POLICE, gras=True, interligne=1.08)

    rect(d, MARGE, 4.92, 1.7, 0.05, fond=OR)

    tf = zone(d, MARGE, 5.16, 10.2, 0.5)
    ecrire(tf, "L'application « RentImmo » — de la formule financière au geste métier",
           taille=15, couleur=TEAL, premier=True)

    tf = zone(d, MARGE, 6.28, 11.0, 0.9)
    ecrire(tf, "HOUNKPETOHOU Marius", taille=15, couleur=BLANC, gras=True,
           premier=True)
    ecrire(tf, [("Choubel Consulting, Casablanca", {}),
                ("     ·     ", {"couleur": OR}),
                ("13 juillet — 21 août 2026", {}),
                ("     ·     ", {"couleur": OR}),
                ("Tuteur : ", {}),
                ("à compléter", {"couleur": OR, "gras": True})],
           taille=11.5, couleur=RGBColor(0xB9, 0xC6, 0xD6), espace_avant=3)


def diapo_02_cabinet(prez):
    d = diapo_standard(prez, 2, "Le cabinet et son métier",
                       "Choubel Consulting — conseil en investissement immobilier, Casablanca")
    g = 7.0                                   # largeur de la colonne de gauche

    tf = zone(d, MARGE, HAUT, g - 0.4, 1.2)
    etiquette(tf, "Ce que le cabinet vend", couleur=BLEU, premier=True)
    ecrire(tf, [("Non pas un bien, mais un ", {}),
                ("accompagnement", {"gras": True}),
                (" : transformer un besoin exprimé en une acquisition chiffrée, "
                 "et conduire le dossier jusqu'à la livraison.", {})],
           taille=13.5, interligne=1.22)

    tf = zone(d, MARGE, HAUT + 1.28, g - 0.4, 1.9)
    etiquette(tf, "Ce qu'il traite", couleur=BLEU, premier=True)
    puces(tf, ["4 familles de biens — terrains, villas, appartements, immeubles",
               "5 niveaux de standing — de l'économique au luxe",
               "3 états — neuf, ancien, ou vendu sur plan (VEFA)"], taille=13)

    tf = zone(d, MARGE, HAUT + 3.15, g - 0.4, 1.5)
    etiquette(tf, "Ses clients", couleur=BLEU, premier=True)
    ecrire(tf, [("Investisseurs particuliers, résidents et non-résidents : la ", {}),
                ("nationalité", {"gras": True}), (" et la ", {}),
                ("situation professionnelle", {"gras": True}),
                (" conditionnent l'accès au crédit local. Ce sont donc des "
                 "données de premier entretien, pas des mentions administratives.", {})],
           taille=13.5, interligne=1.22)
    a_completer(tf, "création, effectif, forme juridique du cabinet")

    x = MARGE + g
    largeur = LARG - g
    image_ajustee(d, IMG / "photos/casablanca_maarif.jpg", x, HAUT, largeur, 2.55)
    tf = zone(d, x, HAUT + 2.66, largeur, 0.9)
    ecrire(tf, "Le quartier Maârif, à Casablanca — l'une des zones de recherche "
               "des dossiers traités.", taille=10.5, couleur=ARDOISE,
           interligne=1.12, premier=True)
    ecrire(tf, "Photo karel291, Wikimedia Commons, CC BY 3.0", taille=9,
           couleur=ARDOISE, italique=True, espace_avant=3)

    c = carte(d, x, HAUT + 3.72, largeur, 0.95, bord_gauche=TEAL)
    tf = zone(d, x + 0.30, HAUT + 3.92, largeur - 0.55, 0.6)
    ecrire(tf, "Cette amplitude n'est pas un argument commercial : c'est une "
               "contrainte de conception.", taille=11.5, couleur=NAVY,
           italique=True, interligne=1.15, premier=True)


def diapo_03_fonctionnement(prez):
    d = diapo_standard(prez, 3, "Le fonctionnement observé",
                       "Un processus stable en six étapes — mais écrit nulle part")

    chevrons(d, MARGE, HAUT, LARG, 0.82, [
        ("1", "Recueil\ndu besoin"), ("2", "Recherche"),
        ("3", "Présentation\net visites"), ("4", "Compromis"),
        ("5", "Acte\nnotarié"), ("6", "Livraison")])

    y = HAUT + 1.12
    largeur = (LARG - 2 * 0.28) / 3
    colonnes = [
        ("Outils supports", ["Tableurs construits au cas par cas",
                             "Notes et échanges informels",
                             "Aucun format de restitution commun"]),
        ("Livrables au client", ["Une sélection de biens",
                                 "Un argumentaire chiffré, oral",
                                 "Rien de systématiquement remis"]),
        ("Acteurs du dossier", ["Le conseiller, du début à la fin",
                                "Agences, promoteurs, notaire",
                                "Banque pour le financement"]),
    ]
    for i, (titre, elements) in enumerate(colonnes):
        x = MARGE + i * (largeur + 0.28)
        carte(d, x, y, largeur, 2.05, bord_gauche=BLEU)
        tf = zone(d, x + 0.30, y + 0.24, largeur - 0.55, 1.7)
        etiquette(tf, titre, couleur=BLEU, taille=10.5, premier=True)
        puces(tf, elements, taille=11.5, espace=5, interligne=1.1)

    c = carte(d, MARGE, y + 2.30, LARG, 0.92, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.34, y + 2.48, LARG - 0.7, 0.6)
    ecrire(tf, [("Observation.  ", {"gras": True, "couleur": OR}),
                ("Le savoir-faire est ", {}),
                ("dans les têtes, pas dans les procédures", {"gras": True}),
                (". Cela fonctionne tant que les effectifs sont stables — et "
                 "rend coûteuse toute reprise de dossier.", {})],
           taille=12.5, couleur=BLANC, interligne=1.15, premier=True)


def diapo_04_diagnostic(prez):
    d = diapo_standard(prez, 4, "Diagnostic",
                       "Trois constats, une problématique")
    g = 7.55

    tf = zone(d, MARGE, HAUT, g - 0.45, 0.35)
    etiquette(tf, "Ce qui coince", couleur=BLEU, premier=True)

    constats = [
        ("Traçabilité", "un chiffre calculé au cas par cas ne se retrouve pas "
                        "six mois plus tard, ni par un autre conseiller."),
        ("Crédibilité en rendez-vous", "un tableur montré à l'écran ne vaut pas "
                                       "un document remis, chiffré et daté."),
        ("Diversité des opérations", "appartement loué, terrain porté dix ans, "
                                     "immeuble construit puis revendu : trois "
                                     "logiques économiques, un seul outil."),
    ]
    y = HAUT + 0.44
    for i, (titre, texte) in enumerate(constats):
        carte(d, MARGE, y, g - 0.45, 1.14, bord_gauche=TEAL)
        tf = zone(d, MARGE + 0.32, y + 0.20, g - 1.05, 0.85)
        ecrire(tf, [(titre + "  ", {"gras": True, "couleur": BLEU}),
                    (texte, {})], taille=12.5, interligne=1.16, premier=True)
        y += 1.28

    tf = zone(d, MARGE, y + 0.06, g - 0.45, 0.7)
    etiquette(tf, "Gisements d'amélioration", couleur=BLEU, premier=True)
    ecrire(tf, "Homogénéiser le calcul  ·  écrire les conventions  ·  produire "
               "un document client  ·  suivre l'avancement des dossiers.",
           taille=11.5, couleur=ARDOISE, interligne=1.15)

    x = MARGE + g
    largeur = LARG - g
    carte(d, x, HAUT, largeur, 2.5, fond=NAVY, contour=None)
    tf = zone(d, x + 0.34, HAUT + 0.30, largeur - 0.68, 2.0)
    etiquette(tf, "Problématique", couleur=OR, premier=True)
    ecrire(tf, [("Comment outiller un conseil qui repose sur le ", {}),
                ("jugement du client", {"gras": True, "couleur": OR}),
                (", sans le remplacer par un verdict automatique ?", {})],
           taille=16, couleur=BLANC, police=TITRE_POLICE, interligne=1.28,
           espace_avant=6)

    y2 = HAUT + 2.85
    for i, (valeur, libelle) in enumerate([("6", "semaines de mission"),
                                           ("2", "développeurs"),
                                           ("8", "questions au tuteur")]):
        chiffre_cle(d, x + i * (largeur / 3), y2, largeur / 3 - 0.1, valeur, libelle)


def diapo_05_intervention(prez):
    d = diapo_standard(prez, 5, "Mon périmètre d'intervention",
                       "Ce que j'ai proposé, et comment le travail a été organisé")
    g = 6.35

    tf = zone(d, MARGE, HAUT, g - 0.4, 1.5)
    etiquette(tf, "La proposition", couleur=BLEU, premier=True)
    ecrire(tf, [("RentImmo", {"gras": True}),
                (" : une application web interne au cabinet, du recueil du besoin "
                 "au document remis au client, en moins de ", {}),
                ("dix minutes", {"gras": True}),
                (" — la durée d'une séquence de rendez-vous.", {})],
           taille=13.5, interligne=1.22)

    y = HAUT + 1.55
    carte(d, MARGE, y, g - 0.4, 1.35, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.32, y + 0.22, g - 1.0, 1.0)
    etiquette(tf, "La règle posée en semaine 1", couleur=OR, taille=9.5, premier=True)
    ecrire(tf, "« Aucune ligne de logique financière ne sera écrite avant que "
               "les conventions de calcul ne soient rédigées. »",
           taille=13, couleur=BLANC, italique=True, interligne=1.2, espace_avant=4)

    tf = zone(d, MARGE, y + 1.55, g - 0.4, 1.6)
    ecrire(tf, [("Les arbitrages difficiles n'étaient pas du code, mais des ", {}),
                ("définitions", {"gras": True}),
                (" : que compte-t-on dans les charges ? un champ laissé vide "
                 "vaut-il zéro, ou « hériter de la valeur par défaut » ?", {})],
           taille=12.5, couleur=ARDOISE, interligne=1.2, premier=True)

    x = MARGE + g
    largeur = LARG - g
    tf = zone(d, x, HAUT, largeur, 0.35)
    etiquette(tf, "L'organisation, semaine par semaine", couleur=BLEU, premier=True)
    tableau(d, x, HAUT + 0.38, largeur, [
        ["Sem.", "Objet", "Livré"],
        ["S1", "Cadrage", "Cahier des charges, conventions, modèle de données"],
        ["S2", "Moteur et socle", "Moteur financier en Python pur, socle web"],
        ["S3", "Restitution", "Graphiques, comparaison de scénarios"],
        ["S4", "Documents client", "Travaux par postes, exports PDF et Excel"],
        ["S5", "Consolidation", "Cas réels, cas limites, jeu de démonstration"],
        ["**S6", "**Cadrage métier", "**Entretien, refonte, étude automatique, mise en ligne"],
    ], largeurs=[0.9, 2.3, 5.4], hauteur_ligne=0.44, taille=11)

    tf = zone(d, x, HAUT + 3.75, largeur, 0.6)
    ecrire(tf, "Un rapport d'avancement chaque dimanche  ·  un jeu de tests au "
               "vert à chaque fin de semaine  ·  tout versionné dans Git.",
           taille=11, couleur=ARDOISE, interligne=1.15, premier=True)


def diapo_06_architecture(prez):
    d = diapo_standard(prez, 6, "Comment c'est construit",
                       "Un choix structurant : isoler le calcul de tout le reste")

    def bloc(x, y, cx, cy, titre, detail, note=None, coeur=False):
        f = rect(d, x, y, cx, cy, fond=RGBColor(0xE6, 0xF7, 0xF5) if coeur else BRUME,
                 contour=TEAL if coeur else TRAIT, epaisseur=1.4 if coeur else 0.75)
        tf = f.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ecrire(tf, titre, taille=12.5, couleur=NAVY, gras=True,
               alignement=PP_ALIGN.CENTER, interligne=1.05, premier=True)
        ecrire(tf, detail, taille=10.5, couleur=ARDOISE,
               alignement=PP_ALIGN.CENTER, interligne=1.05, espace_avant=2)
        if note:
            ecrire(tf, note, taille=9.5, couleur=TEAL, gras=True,
                   alignement=PP_ALIGN.CENTER, interligne=1.05, espace_avant=2)
        return f

    def fleche(x1, y1, x2, y2):
        from pptx.enum.shapes import MSO_CONNECTOR
        c = d.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
        c.line.color.rgb = BLEU
        c.line.width = Pt(1.5)
        _pointe_de_fleche(c)

    y1, hb = HAUT + 0.30, 1.05
    bloc(MARGE, y1, 2.5, hb, "Navigateur", "Bootstrap, Chart.js")
    bloc(MARGE + 3.05, y1, 2.9, hb, "Application Flask", "pages, formulaires, sécurité")
    bloc(MARGE + 6.50, y1 - 0.15, 3.15, hb + 0.30, "Moteur financier",
         "Python pur, sans Flask", "testable seul", coeur=True)
    y2 = y1 + 1.75
    bloc(MARGE + 3.05, y2, 2.9, hb, "Base de données", "SQLAlchemy + migrations")
    bloc(MARGE + 6.50, y2, 3.15, hb, "Exports", "PDF (WeasyPrint), Excel")

    fleche(MARGE + 2.50, y1 + hb / 2, MARGE + 3.05, y1 + hb / 2)
    fleche(MARGE + 5.95, y1 + hb / 2, MARGE + 6.50, y1 + hb / 2)
    fleche(MARGE + 4.50, y1 + hb, MARGE + 4.50, y2)
    fleche(MARGE + 8.07, y1 + hb + 0.15, MARGE + 8.07, y2)

    y3 = y2 + 1.42
    largeur = (LARG - 0.34) / 2
    carte(d, MARGE, y3, largeur, 1.42, bord_gauche=BLEU)
    tf = zone(d, MARGE + 0.32, y3 + 0.22, largeur - 0.6, 1.05)
    etiquette(tf, "Pourquoi séparer le moteur ?", couleur=BLEU, taille=10, premier=True)
    puces(tf, ["il se teste sans lancer l'application",
               "il se confronte à des valeurs de référence externes",
               "il est réutilisable ailleurs"], taille=11.5, espace=3, interligne=1.08)

    x = MARGE + largeur + 0.34
    carte(d, x, y3, largeur, 1.42, bord_gauche=TEAL)
    tf = zone(d, x + 0.32, y3 + 0.22, largeur - 0.6, 1.05)
    etiquette(tf, "Une décision de conception", couleur=BLEU, taille=10, premier=True)
    ecrire(tf, [("Les résultats ne sont ", {}),
                ("jamais stockés", {"gras": True}),
                (", toujours recalculés : une incohérence entre les hypothèses "
                 "affichées et les indicateurs devient structurellement "
                 "impossible.", {})],
           taille=11.5, interligne=1.14, espace_avant=3)


def _pointe_de_fleche(connecteur):
    """Ajoute une pointe de flèche à l'extrémité d'un connecteur."""
    from pptx.oxml.ns import qn
    ln = connecteur.line._get_or_add_ln()
    fin = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(fin)


def diapo_07_brief(prez):
    d = diapo_standard(prez, 7, "Réalisation — partir du besoin, pas du bien",
                       "Le brief de recherche : les critères recueillis au premier entretien")
    image_ajustee(d, IMG / "slides/app_brief.png", MARGE, HAUT, LARG, 3.72)

    y = HAUT + 3.88
    largeur = (LARG - 0.34) / 2
    carte(d, MARGE, y, largeur, 1.14, bord_gauche=BLEU)
    tf = zone(d, MARGE + 0.32, y + 0.20, largeur - 0.6, 0.85)
    ecrire(tf, "Type de bien, standing, superficie, distribution, commodités, "
               "zone, budget, mode de financement — ajustés au bien cherché.",
           taille=12.5,
           interligne=1.16, premier=True)

    x = MARGE + largeur + 0.34
    carte(d, x, y, largeur, 1.14, fond=NAVY, contour=None)
    tf = zone(d, x + 0.32, y + 0.20, largeur - 0.6, 0.85)
    ecrire(tf, [("Et surtout : l'", {}),
                ("objectif", {"gras": True, "couleur": OR}), (" et l'", {}),
                ("horizon", {"gras": True, "couleur": OR}),
                (" du client, qui serviront ensuite à lire les indicateurs. "
                 "Cet écran n'existait pas avant la semaine 6.", {})],
           taille=12.5, couleur=BLANC, interligne=1.16, premier=True)


def diapo_08_restitution(prez):
    d = diapo_standard(prez, 8, "Réalisation — ce que le conseiller montre au client",
                       "Les indicateurs, puis les montages comparés côte à côte")
    largeur = (LARG - 0.4) / 2

    for i, (fichier, texte) in enumerate([
        ("slides/app_resultats_locatif.png",
         [("Rendement net", {"gras": True}), (" et ", {}),
          ("cash-flow", {"gras": True}),
          (" en première ligne. Valeur créée, TRI et VAN au second rang. ", {}),
          ("Aucun seuil de rentabilité", {"gras": True, "couleur": BLEU}),
          (" : ils se lisent à l'aune de l'objectif du client.", {})]),
        ("slides/app_comparaison.png",
         [("Trois montages pour le même bien : l'écran d'arbitrage du "
           "rendez-vous. Chaque scénario s'exporte en ", {}),
          ("PDF", {"gras": True}), (" et en ", {}), ("Excel", {"gras": True}),
          (" aux couleurs du cabinet.", {})]),
    ]):
        x = MARGE + i * (largeur + 0.4)
        image_ajustee(d, IMG / fichier, x, HAUT, largeur, 2.35)
        carte(d, x, HAUT + 2.52, largeur, 1.05, bord_gauche=TEAL)
        tf = zone(d, x + 0.30, HAUT + 2.70, largeur - 0.58, 0.8)
        ecrire(tf, texte, taille=12, interligne=1.16, premier=True)

    c = carte(d, MARGE, HAUT + 3.80, LARG, 1.02, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.34, HAUT + 4.00, LARG - 0.7, 0.7)
    ecrire(tf, [("Les conventions de calcul sont imprimées dans chaque export", {"gras": True}),
                (" — périmètre des charges, fiscalité au taux effectif, travaux "
                 "intégrés au coût d'acquisition. Le conseiller ne peut pas être "
                 "contredit sur un chiffre dont la convention est écrite sous le "
                 "tableau.", {})],
           taille=12.5, couleur=BLANC, interligne=1.16, premier=True)


def diapo_09_tournant(prez):
    d = diapo_standard(prez, 9, "Le tournant du stage",
                       "Semaine 6 : huit questions au tuteur — et trois hypothèses sur quatre invalidées")
    g = 6.6

    carte(d, MARGE, HAUT, g - 0.4, 3.05, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.40, HAUT + 0.30, g - 1.2, 2.5)
    ecrire(tf, "“", taille=40, couleur=OR, police=TITRE_POLICE, gras=True,
           interligne=0.6, premier=True)
    ecrire(tf, [("Les deux critères les plus récurrents sont le ", {}),
                ("rendement net", {"gras": True, "couleur": OR}), (" et le ", {}),
                ("cash-flow", {"gras": True, "couleur": OR}), (".", {})],
           taille=13, couleur=BLANC, interligne=1.2, espace_apres=8)
    ecrire(tf, [("Un investissement est bon lorsqu'il permet au client de ", {}),
                ("générer de la valeur", {"gras": True, "couleur": OR}),
                (". Tout dépend donc de son ", {}),
                ("horizon", {"gras": True, "couleur": OR}), (" et du ", {}),
                ("facteur temps", {"gras": True, "couleur": OR}), (".", {})],
           taille=13, couleur=BLANC, interligne=1.2, espace_apres=8)
    ecrire(tf, "Entretien de cadrage métier, août 2026", taille=10,
           couleur=RGBColor(0x8E, 0x9E, 0xB4), italique=True)

    x = MARGE + g
    largeur = LARG - g
    for i, (titre, couleur, elements) in enumerate([
        ("Ce que je croyais", ARDOISE,
         ["Le TRI et la VAN décident.",
          "Il existe un seuil de rentabilité.",
          "Tout bien produit un loyer."]),
        ("Ce que le métier fait", TEAL,
         ["Rendement net et cash-flow d'abord.",
          "Aucun seuil : le client tranche.",
          "Les deux meilleurs dossiers n'ont aucun loyer."]),
    ]):
        y = HAUT + i * 1.62
        carte(d, x, y, largeur, 1.42, bord_gauche=couleur)
        tf = zone(d, x + 0.32, y + 0.20, largeur - 0.6, 1.1)
        etiquette(tf, titre, couleur=couleur, taille=10, premier=True)
        puces(tf, elements, taille=11.5, puce_couleur=couleur, espace=3,
              interligne=1.08)

    c = carte(d, MARGE, HAUT + 3.30, LARG, 1.05, bord_gauche=ROUGE)
    tf = zone(d, MARGE + 0.34, HAUT + 3.50, LARG - 0.7, 0.75)
    ecrire(tf, [("Mon erreur de méthode.  ", {"gras": True, "couleur": ROUGE}),
                ("Cet entretien aurait dû avoir lieu en semaine 1, pas en "
                 "semaine 6. Cinq semaines ont reposé sur des hypothèses non "
                 "vérifiées : j'avais pris le silence de départ pour un accord.", {})],
           taille=12.5, interligne=1.16, premier=True)


def diapo_10_dossiers(prez):
    d = diapo_standard(prez, 10, "Les deux dossiers qui ont tout changé",
                       "Même capital, même cabinet — le classement dépend de la durée")
    g = 8.55

    tableau(d, MARGE, HAUT, g - 0.4, [
        ["", "Client 1 — portage foncier", "Client 2 — construction-revente"],
        ["Budget", "≈ 1 000 000 MAD", "≈ 1 000 000 MAD"],
        ["Bien", "Terrain de 3 ha, zone ciblée", "Terrain viabilisé de 500 m² à 1,1 M"],
        ["Opération", "Conservé, sans mise en valeur", "Immeuble R+4, 20 appartements revendus"],
        ["**Horizon", "**10 ans", "**2 ans"],
        ["**Résultat", "**≈ 15 000 000 MAD de bénéfice", "**5 000 000 MAD de plus-value"],
    ], largeurs=[1.6, 3.4, 3.6], hauteur_ligne=0.50, hauteur_entete=0.46, taille=11.5)

    y = HAUT + 3.12
    carte(d, MARGE, y, g - 0.4, 1.22, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.34, y + 0.22, g - 1.1, 0.9)
    ecrire(tf, [("À deux ans, le client 2 a fait la meilleure opération. "
                 "À dix ans, le classement s'inverse. ", {}),
                ("Aucun des deux dossiers n'a de loyer.", {"gras": True, "couleur": OR})],
           taille=13, couleur=BLANC, interligne=1.18, premier=True)

    x = MARGE + g
    largeur = LARG - g
    image_ajustee(d, IMG / "photos/casablanca_marina_chantier.jpg", x, HAUT,
                  largeur, 2.6)
    tf = zone(d, x, HAUT + 2.72, largeur, 1.5)
    ecrire(tf, "Un chantier à Casablanca : sur un achat sur plan, les mensualités "
               "courent pendant que le bien ne produit encore rien.",
           taille=11, couleur=ARDOISE, interligne=1.15, premier=True)
    ecrire(tf, "Photo Farid Mernissi, Wikimedia Commons, CC BY-SA 4.0", taille=9,
           couleur=ARDOISE, italique=True, espace_avant=4)


def diapo_11_refonte(prez):
    d = diapo_standard(prez, 11, "Ce que la refonte a changé",
                       "Sept écarts relevés, sept corrections apportées en semaine 6")

    tableau(d, MARGE, HAUT, LARG, [
        ["#", "Écart constaté", "Correction apportée"],
        ["1", "TRI et VAN affichés en tête", "Rendement net et cash-flow remontés en première ligne, partout"],
        ["2", "Un loyer supposé dans tous les cas", "Opération sans loyer traitée nativement : la valeur vient de la plus-value"],
        ["3", "Le facteur temps, paramètre secondaire", "Horizon et délai de livraison passés au premier plan"],
        ["4", "Phase amont absente", "Ajout du brief de recherche au modèle de données"],
        ["5", "Fiche client incomplète", "Situation professionnelle, nationalité, budget disponible"],
        ["6", "Déroulé du dossier non suivi", "Les six étapes du métier deviennent six statuts de dossier"],
        ["7", "VEFA traitée comme un bien louable", "Délai de livraison modélisé, et proratisé sur l'année"],
    ], largeurs=[0.5, 4.3, 7.4], hauteur_ligne=0.42, hauteur_entete=0.44, taille=11.5)

    y = HAUT + 3.44
    carte(d, MARGE, y, LARG, 1.05, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.34, y + 0.22, LARG - 0.7, 0.75)
    ecrire(tf, [("Et une suppression : ", {}),
                ("le seuil de rentabilité universel", {"gras": True, "couleur": OR}),
                (". Un montage n'est plus jugé dans l'absolu, mais relativement "
                 "à l'objectif et à l'horizon que le client a déclarés.", {})],
           taille=13, couleur=BLANC, interligne=1.16, premier=True)


def diapo_12_etude(prez):
    d = diapo_standard(prez, 12, "L'étude automatique",
                       "La dernière fonctionnalité livrée — et celle dont je suis le plus satisfait")

    chevrons(d, MARGE, HAUT, LARG, 0.70, [
        ("1", "Le bien et le client"),
        ("2", "Sept montages construits"),
        ("3", "Confrontés sur quatre critères"),
        ("4", "Un montage proposé, et expliqué")], actif=3)

    y = HAUT + 0.98
    largeur = (LARG - 0.4) / 2
    x = MARGE + largeur + 0.4

    image_ajustee(d, IMG / "slides/app_etude.png", MARGE, y, largeur, 2.36)

    carte(d, x, y, largeur, 2.36, bord_gauche=TEAL)
    tf = zone(d, x + 0.32, y + 0.22, largeur - 0.62, 2.0)
    etiquette(tf, "Proposer sans réintroduire un seuil", couleur=BLEU,
              taille=10, premier=True)
    puces(tf, ["Classement relatif — ramené entre le meilleur et le moins bon "
               "montage de cette étude, pour ce bien et cet horizon",
               "Pondérations affichées — le conseiller voit sur quoi l'outil a "
               "tranché, et peut le contester",
               "Ce que les autres font mieux — nommer ce que l'on perd, sinon "
               "c'est de la vente et non du conseil",
               "Composition complète — chaque valeur avec son origine : dossier, "
               "brief, zone de marché, calcul ou hypothèse"],
          taille=10.5, espace=5, interligne=1.06)

    y2 = y + 2.56
    carte(d, MARGE, y2, largeur, 1.20, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.32, y2 + 0.20, largeur - 0.62, 0.9)
    ecrire(tf, [("Le conseiller saisissait ", {}),
                ("onze valeurs", {"gras": True, "couleur": OR}),
                (" pour découvrir que le montage ne convenait pas, puis "
                 "recommençait : l'outil lui demandait de deviner la réponse "
                 "avant de la calculer.", {})],
           taille=11.5, couleur=BLANC, interligne=1.12, premier=True)

    carte(d, x, y2, largeur, 1.20, bord_gauche=BLEU)
    tf = zone(d, x + 0.32, y2 + 0.20, largeur - 0.62, 0.9)
    ecrire(tf, [("À bien et budget identiques", {"gras": True}),
                (", l'outil propose ", {}),
                ("vingt-cinq ans", {"gras": True, "couleur": BLEU}),
                (" au client qui cherche du revenu, et ", {}),
                ("quinze ans", {"gras": True, "couleur": BLEU}),
                (" à celui qui vise la plus-value. Un test le vérifie.", {})],
           taille=11.5, interligne=1.12, premier=True)


def diapo_13_validation(prez):
    d = diapo_standard(prez, 13, "Tests, validation et passation",
                       "Rendre le travail vérifiable, puis utilisable sans moi")
    largeur = (LARG - 0.4) / 2

    carte(d, MARGE, HAUT, largeur, 3.55, bord_gauche=BLEU)
    tf = zone(d, MARGE + 0.34, HAUT + 0.26, largeur - 0.66, 3.0)
    etiquette(tf, "Comment c'est validé", couleur=BLEU, taille=10.5, premier=True)
    puces(tf, ["96 tests automatisés — quatre niveaux : moteur, parcours "
               "complet, cas réels, règles métier",
               "Valeurs de référence externes — une mensualité vérifiée au centime",
               "Les deux dossiers du tuteur — devenus des tests : 15 M et 5 M de "
               "valeur créée, recalculables à la main",
               "Deux tests-verdict — le classement s'inverse avec l'horizon, et "
               "l'objectif du client change le montage proposé : la position du "
               "cabinet est inscrite dans le code"],
          taille=12, espace=7, interligne=1.12)

    x = MARGE + largeur + 0.4
    carte(d, x, HAUT, largeur, 3.55, bord_gauche=TEAL)
    tf = zone(d, x + 0.34, HAUT + 0.26, largeur - 0.66, 3.0)
    etiquette(tf, "Comment c'est transmis", couleur=BLEU, taille=10.5, premier=True)
    puces(tf, ["Application déployée en ligne — accessible sans installation",
               "Jeu de démonstration — rejouable en une commande",
               "Guide d'utilisation — rédigé dans l'ordre d'un rendez-vous",
               "Script de démonstration — déroulé minuté de dix minutes"],
          taille=12, puce_couleur=TEAL, espace=7, interligne=1.12)
    a_completer(tf, "retour des utilisateurs après prise en main")

    y = HAUT + 3.78
    carte(d, MARGE, y, LARG, 0.92, bord_gauche=ROUGE)
    tf = zone(d, MARGE + 0.34, y + 0.20, LARG - 0.7, 0.6)
    ecrire(tf, [("Un test a échoué parce que ", {}), ("ma", {"gras": True, "italique": True}),
                (" valeur de référence était fausse — le programme, lui, avait "
                 "raison. Le dispositif protège dans les deux sens.", {})],
           taille=12.5, interligne=1.16, premier=True)


def diapo_14_conclusion(prez):
    d = diapo_standard(prez, 14, "Conclusion et perspectives",
                       "Un outil qui propose un montage — et documente la décision au lieu de la remplacer")
    g = 6.9

    tf = zone(d, MARGE, HAUT, g - 0.4, 0.35)
    etiquette(tf, "Ce qui est livré", couleur=BLEU, premier=True)
    tf = zone(d, MARGE, HAUT + 0.42, g - 0.4, 1.7)
    puces(tf, ["Le parcours complet, de la fiche client au document remis",
               "Les trois familles d'opérations, avec ou sans loyer",
               "Le facteur temps au cœur du calcul",
               "L'étude automatique : les montages construits et confrontés",
               "Documentation, tests, déploiement, application en ligne"],
          taille=13, espace=7, interligne=1.14, premier=True)

    tf = zone(d, MARGE, HAUT + 2.24, g - 0.4, 0.35)
    etiquette(tf, "Ce qui reste ouvert", couleur=ARDOISE, premier=True)
    tf = zone(d, MARGE, HAUT + 2.66, g - 0.4, 1.7)
    puces(tf, ["Échéancier VEFA par tranches (appels de fonds)",
               "Calcul de la capacité d'emprunt",
               "Rapprochement entre le brief et un portefeuille de biens",
               "Validation des taux par défaut par le cabinet"],
          taille=13, puce_couleur=ARDOISE, espace=7, interligne=1.14,
          premier=True)

    x = MARGE + g
    largeur = LARG - g
    carte(d, x, HAUT, largeur, 2.75, fond=RGBColor(0xE6, 0xF7, 0xF5),
          contour=TEAL)
    tf = zone(d, x + 0.36, HAUT + 0.32, largeur - 0.72, 2.2)
    ecrire(tf, "Le déplacement le plus solide du stage n'est pas technique :",
           taille=12.5, couleur=ARDOISE, interligne=1.15, premier=True)
    ecrire(tf, "d'un outil qui décide", taille=16, couleur=NAVY,
           police=TITRE_POLICE, gras=True, espace_avant=8, interligne=1.1)
    ecrire(tf, "→  vers un outil qui documente une décision.", taille=16,
           couleur=BLEU, police=TITRE_POLICE, gras=True, espace_avant=2,
           interligne=1.1)
    ecrire(tf, "Il propose — mais il montre la règle qui l'a fait choisir.",
           taille=11.5, couleur=ARDOISE, italique=True, espace_avant=6,
           interligne=1.12)
    ecrire(tf, "Il ne vient pas d'une compétence, mais d'avoir posé des questions "
               "à quelqu'un qui exerce le métier.", taille=11.5, couleur=ARDOISE,
           espace_avant=10, interligne=1.15)

    y = HAUT + 3.00
    for i, (valeur, libelle) in enumerate([("96", "tests au vert"),
                                           ("6", "étapes suivies"),
                                           ("7", "écarts corrigés")]):
        chiffre_cle(d, x + i * (largeur / 3), y, largeur / 3 - 0.1, valeur, libelle)

    y = HAUT + 4.20
    carte(d, MARGE, y, LARG, 1.00, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.36, y + 0.21, LARG - 0.72, 0.70)
    ecrire(tf, [("La règle du cabinet, désormais inscrite dans l'outil.  ",
                 {"gras": True, "couleur": OR}),
                ("« Un investissement est bon lorsqu'il permet au client de "
                 "générer de la valeur. Tout dépend donc de son horizon "
                 "d'investissement et du facteur temps. »", {"italique": True})],
           taille=12.5, couleur=BLANC, interligne=1.16, premier=True)


def diapo_15_bilan(prez):
    d = diapo_standard(prez, 15, "Bilan personnel",
                       "Ce que j'ai appris, ce qui m'a manqué, et ce que j'en fais")
    largeur = (LARG - 2 * 0.30) / 3

    colonnes = [
        ("Ce que la 1A a servi", BLEU,
         [("Méthodes numériques", "chercher le TRI, c'est résoudre une équation : "
                                  "Newton-Raphson, puis repli sur une bissection bornée."),
          ("Mathématiques financières", "actualisation, annuités constantes."),
          ("Programmation et communication écrite", "en continu.")]),
        ("Ce que j'ai appris", TEAL,
         [("Écrire les conventions avant le code", "les arbitrages durs sont des "
                                                   "définitions, pas des algorithmes."),
          ("Vérifier deux fois", "le dispositif protège aussi contre mes propres erreurs."),
          ("Si je ne sais pas l'expliquer", "je ne dois pas le livrer.")]),
        ("Ce qui m'a manqué", ROUGE,
         [("Conduire un entretien utilisateur", "mon questionnaire a été improvisé, "
                                                "et trop tardif."),
          ("Culture du marché immobilier", "découverte en même temps que le vocabulaire."),
          ("Réflexes d'exploitation", "un bug invisible en local m'a appris la "
                                      "différence entre « ça marche » et « ça tourne ».")]),
    ]
    for i, (titre, couleur, elements) in enumerate(colonnes):
        x = MARGE + i * (largeur + 0.30)
        carte(d, x, HAUT, largeur, 3.62, bord_gauche=couleur)
        tf = zone(d, x + 0.30, HAUT + 0.26, largeur - 0.58, 3.2)
        etiquette(tf, titre, couleur=couleur, taille=10, premier=True)
        for tete, reste in elements:
            ecrire(tf, [(tete + " ", {"gras": True}), ("— " + reste, {})],
                   taille=12.5, interligne=1.16, espace_apres=10)
        if i == 0:
            a_completer(tf, "intitulés des cours")

    y = HAUT + 3.86
    carte(d, MARGE, y, LARG, 1.32, fond=NAVY, contour=None)
    tf = zone(d, MARGE + 0.36, y + 0.28, LARG - 0.72, 0.95)
    ecrire(tf, [("Mon plan.  ", {"gras": True, "couleur": OR}),
                ("Imposer un ", {}),
                ("entretien de cadrage écrit dès le premier jour", {"gras": True}),
                (" de ma prochaine mission  ·  pratiquer le déploiement et la "
                 "supervision sur mes propres projets  ·  chercher un stage à "
                 "l'intersection du logiciel et d'un ", {}),
                ("métier exigeant", {"gras": True}), (", dans une ", {}),
                ("équipe étoffée", {"gras": True}),
                (" — ce stage s'est fait en binôme, avec une autonomie large.", {})],
           taille=12.5, couleur=BLANC, interligne=1.18, premier=True)


def diapo_16_merci(prez):
    d = diapo_photo(prez, IMG / "photos/casablanca_vue_aerienne.jpg", 0.88)

    tf = zone(d, MARGE, 2.60, 11.0, 0.4)
    etiquette(tf, "Choubel Consulting — 13 juillet au 21 août 2026", premier=True)

    tf = zone(d, MARGE, 3.05, 11.0, 1.1)
    ecrire(tf, "Merci de votre attention", taille=40, couleur=BLANC,
           police=TITRE_POLICE, gras=True, interligne=1.05, premier=True)

    rect(d, MARGE, 4.32, 1.7, 0.05, fond=OR)

    tf = zone(d, MARGE, 4.58, 10.0, 0.4)
    ecrire(tf, "Questions bienvenues", taille=16, couleur=TEAL, premier=True)

    tf = zone(d, MARGE, 6.05, 11.6, 1.0)
    ecrire(tf, "Crédits photographiques", taille=10, couleur=OR, gras=True,
           premier=True)
    ecrire(tf, "Casablanca Finance City : Turdyfirst, CC BY-SA 4.0  ·  quartier "
               "Maârif : karel291, CC BY 3.0  ·  marina en construction : Farid "
               "Mernissi, CC BY-SA 4.0  ·  vue aérienne : Jimmy Baikovicius, "
               "CC BY-SA 2.0 — toutes via Wikimedia Commons. Les captures d'écran "
               "sont des productions du stage.",
           taille=9.5, couleur=RGBColor(0x9A, 0xAA, 0xBE), interligne=1.2,
           espace_avant=3)


# ══════════════════════════════════════════════════════════════════════════════
def construire():
    prez = Presentation()
    prez.slide_width, prez.slide_height = Inches(L), Inches(H)

    for fabriquer in (diapo_01_titre, diapo_02_cabinet, diapo_03_fonctionnement,
                      diapo_04_diagnostic, diapo_05_intervention,
                      diapo_06_architecture, diapo_07_brief,
                      diapo_08_restitution, diapo_09_tournant,
                      diapo_10_dossiers, diapo_11_refonte, diapo_12_etude,
                      diapo_13_validation, diapo_14_conclusion,
                      diapo_15_bilan, diapo_16_merci):
        fabriquer(prez)

    # Le pied de page annonce un total : qu'il cesse d'être vrai si l'on ajoute
    # une diapositive sans y penser serait une faute discrète et durable.
    assert len(prez.slides._sldIdLst) == NB_DIAPOS, (
        f"NB_DIAPOS vaut {NB_DIAPOS} mais la présentation en compte "
        f"{len(prez.slides._sldIdLst)}"
    )

    prez.core_properties.title = "Concevoir un outil d'aide à la décision pour un cabinet de conseil immobilier"
    prez.core_properties.author = "HOUNKPETOHOU Marius"
    prez.core_properties.subject = "Soutenance de stage opérateur — École Centrale Casablanca"
    prez.save(SORTIE)
    print(f"{SORTIE.name} — {len(prez.slides._sldIdLst)} diapositives")


if __name__ == "__main__":
    construire()
